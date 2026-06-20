"""Module for handling various input types and converting them to Markdown.

Supports text, documents, spreadsheets, presentations, images (OCR),
PDFs (with table detection), HTML, audio, video, and YouTube transcripts.

Heavy optional dependencies (PyMuPDF, Tesseract/Pillow, faster-whisper,
youtube-transcript-api, yt-dlp) are imported lazily so that the core package
stays importable and lightweight without the corresponding extras.
"""

from __future__ import annotations

import importlib
import os
import re
import subprocess
import tempfile
import threading
import warnings
from html.parser import HTMLParser
from pathlib import Path
from typing import TYPE_CHECKING, Any, Dict, List, Optional, Tuple
from uuid import uuid4

import pandas as pd
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation

if TYPE_CHECKING:
    import fitz
    from faster_whisper import WhisperModel

# PyMuPDF span flag bits: bit 1 (value 2) is italic, bit 4 (value 16) is bold.
_BOLD_FLAG: int = 1 << 4

# Default Whisper model size; can be overridden per call or via the
# ANY_TO_MARKDOWN_WHISPER_MODEL environment variable.
_DEFAULT_WHISPER_MODEL: str = "small"
_WHISPER_MODEL_ENV_VAR: str = "ANY_TO_MARKDOWN_WHISPER_MODEL"

# Cache of Faster-Whisper model instances keyed by model size, so switching
# sizes never silently reuses the wrong model. Initialized lazily to save
# resources if no audio/video files are processed.
_whisper_models: Dict[str, "WhisperModel"] = {}
_model_lock = threading.Lock()


class MissingDependencyError(ImportError):
    """Raised when an optional dependency required by a handler is not installed."""


class TranscriptUnavailableError(RuntimeError):
    """Raised when a YouTube transcript cannot be retrieved."""


def require_dependency(module_name: str, extra: str) -> Any:
    """Imports an optional dependency or raises an actionable error.

    Args:
        module_name: The importable module name (e.g. 'fitz', 'PIL.Image').
        extra: The pip extra that provides the module (e.g. 'pdf').

    Returns:
        The imported module.

    Raises:
        MissingDependencyError: If the module is not installed, with the exact
            pip command needed to fix it.
    """
    try:
        return importlib.import_module(module_name)
    except ImportError as e:
        raise MissingDependencyError(
            f"The '{module_name}' package is required for this feature. "
            f"Install it with: pip install any-to-markdown[{extra}]"
        ) from e


def get_whisper_model(model_size: Optional[str] = None) -> WhisperModel:
    """Lazily initializes and returns a Faster-Whisper model.

    Architectural Decision:
    - Lazy initialization ensures we don't consume memory/GPU resources until needed.
    - Instances are cached per model size and guarded by a global lock, so
      switching sizes never reuses the wrong model.
    - Uses GPU (CUDA) if available, otherwise falls back to CPU with int8 quantization.

    Args:
        model_size: Whisper model size (e.g. 'tiny', 'small', 'medium',
            'large-v3'). Defaults to the ANY_TO_MARKDOWN_WHISPER_MODEL
            environment variable, or 'small' if unset.

    Returns:
        WhisperModel: An initialized instance of the Faster-Whisper model.

    Raises:
        MissingDependencyError: If the 'audio' extra is not installed.
    """
    if model_size is None:
        model_size = os.environ.get(_WHISPER_MODEL_ENV_VAR, _DEFAULT_WHISPER_MODEL)

    with _model_lock:
        if model_size not in _whisper_models:
            ctranslate2 = require_dependency("ctranslate2", "audio")
            faster_whisper = require_dependency("faster_whisper", "audio")

            # Auto-detect CUDA capability
            device = "cuda" if ctranslate2.get_cuda_device_count() > 0 else "cpu"

            # Optimization: float16 is significantly faster on modern GPUs,
            # while int8 is highly optimized for modern CPUs via ctranslate2.
            compute_type = "float16" if device == "cuda" else "int8"

            _whisper_models[model_size] = faster_whisper.WhisperModel(
                model_size, device=device, compute_type=compute_type
            )
        return _whisper_models[model_size]


def handle_text(file_path: str | Path) -> str:
    """Reads plain text files with tolerant decoding and returns content.

    Decoding strategy: UTF-8 (BOM-aware via utf-8-sig) first, then Latin-1 as
    a fallback, so a single odd byte never fails the whole file. Latin-1 maps
    every byte to a code point, making the fallback lossless and total.

    Args:
        file_path: Path to the source text file.

    Returns:
        The file content wrapped in Markdown-friendly spacing.
    """
    raw = Path(file_path).read_bytes()
    try:
        user_data = raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        user_data = raw.decode("latin-1")
    return f"\n\n{user_data}\n\n"


def handle_document(file_path: str | Path) -> str:
    """Processes DOCX files, preserving the logical order of paragraphs and tables.

    Implementation Detail:
    - Iterates through the document's XML body directly to ensure that
      tables are interleaved with text correctly.

    Args:
        file_path: Path to the .docx file.

    Returns:
        Consolidated Markdown content of the document.
    """
    user_doc = Document(str(file_path))
    parts: List[str] = []

    # Direct access to the body elements allows us to maintain strict document order
    for element in user_doc.element.body:
        if isinstance(element, CT_P):
            para = Paragraph(element, user_doc)
            if para.text.strip():
                parts.append(para.text + "\n\n")
        elif isinstance(element, CT_Tbl):
            table = Table(element, user_doc)
            table_data: List[List[str]] = []
            for row in table.rows:
                table_data.append([cell.text.strip() for cell in row.cells])

            if table_data:
                df = pd.DataFrame(table_data)
                # Heuristic: If multiple rows exist, assume the first row is a header
                if len(table_data) > 1:
                    df.columns = pd.Index(table_data[0])
                    df = df.iloc[1:]
                parts.append(df.to_markdown(index=False) + "\n\n")

    return "".join(parts)


def handle_excel(file_path: str | Path) -> str:
    """Processes Excel files (.xls, .xlsx).

    Args:
        file_path: Path to the Excel spreadsheet.

    Returns:
        All sheets converted into Markdown tables.
    """
    # Load all sheets to ensure no data is missed
    user_excel: Dict[str, pd.DataFrame] = pd.read_excel(file_path, sheet_name=None)
    parts: List[str] = []
    for sheet_name, df in user_excel.items():
        parts.append(f"### Sheet: {sheet_name}\n\n")
        parts.append(f"{df.to_markdown(index=False)}\n\n")
    return "".join(parts)


def handle_csv(file_path: str | Path) -> str:
    """Processes CSV files into a Markdown table.

    Args:
        file_path: Path to the .csv file.

    Returns:
        The CSV content converted into a Markdown table.
    """
    df = pd.read_csv(file_path)
    return f"{df.to_markdown(index=False)}\n\n"


def handle_powerpoint(file_path: str | Path) -> str:
    """Processes PowerPoint files, extracting text slide by slide.

    Extracts titles, text boxes, tables, and speaker notes.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        Segmented Markdown content per slide.
    """
    ppt = Presentation(str(file_path))
    parts: List[str] = []
    for i, slide in enumerate(ppt.slides, start=1):
        parts.append(f"## Slide {i}\n\n")

        # Extract slide title as a level 3 header
        title = slide.shapes.title
        if title and hasattr(title, "text") and title.text.strip():
            parts.append(f"### {title.text.strip()}\n\n")

        # Process all shapes (text boxes, tables, etc.)
        for shape in slide.shapes:
            if shape == title:
                continue

            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip() + "\n\n")

            if shape.has_table:
                table_data: List[List[str]] = []
                for row in shape.table.rows:
                    table_data.append([cell.text_frame.text.strip() for cell in row.cells])
                df = pd.DataFrame(table_data)
                parts.append(df.to_markdown(index=False) + "\n\n")

        # Capture speaker notes to provide context often missed in raw slides
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"*Speaker Notes:* {notes}\n\n")

    return "".join(parts)


def handle_image(file_path: str | Path) -> str:
    """Uses Tesseract OCR to extract text from images.

    Args:
        file_path: Path to the image file.

    Returns:
        The extracted text content.

    Raises:
        MissingDependencyError: If the 'ocr' extra is not installed.
    """
    pytesseract = require_dependency("pytesseract", "ocr")
    pil_image = require_dependency("PIL.Image", "ocr")

    with pil_image.open(file_path) as image:
        # Pre-processing for better OCR accuracy
        grayscale = image.convert("L")
        text: str = pytesseract.image_to_string(grayscale, config="--psm 6")
    return text.strip()


def _get_pdf_items(page: fitz.Page, ignore_bboxes: Optional[List[fitz.Rect]] = None) -> List[Tuple[float, str]]:
    """Internal helper for PDF extraction that preserves reading order.

    Args:
        page: The PyMuPDF page object.
        ignore_bboxes: Regions to exclude (e.g., tables).

    Returns:
        List of (y_coordinate, content) tuples.
    """
    fitz_mod = require_dependency("fitz", "pdf")

    if ignore_bboxes is None:
        ignore_bboxes = []

    blocks = page.get_text("dict")["blocks"]
    items: List[Tuple[float, str]] = []

    for b in blocks:
        if b["type"] == 0:  # text block
            # Spatial exclusion: Skip text that overlaps with detected table regions
            if any(fitz_mod.Rect(b["bbox"]).intersects(ib) for ib in ignore_bboxes):
                continue

            y = b["bbox"][1]
            structured_parts: List[str] = []
            for l in b["lines"]:  # noqa: E741
                for s in l["spans"]:
                    text = s["text"].strip()
                    if not text:
                        continue

                    size = s["size"]
                    flags = s["flags"]

                    # Heuristic for Markdown conversion based on PDF styling
                    if size > 16:
                        structured_parts.append(f"\n# {text}\n")
                    elif size > 13:
                        structured_parts.append(f"\n## {text}\n")
                    elif flags & _BOLD_FLAG:  # bold flag bit in PyMuPDF (bit 4 = 16)
                        structured_parts.append(f"\n### {text}\n")
                    else:
                        structured_parts.append(text + " ")
                structured_parts.append("\n")
            items.append((y, "".join(structured_parts)))
    return items


def handle_pdf(file_path: str | Path, use_layout_engine: bool = False) -> str:
    """Advanced PDF handler using PyMuPDF and OCR fallback.

    Args:
        file_path: Path to the .pdf file.
        use_layout_engine: Whether to use advanced PDF layout analysis.

    Returns:
        The structured Markdown content of the PDF.

    Raises:
        MissingDependencyError: If the 'pdf' extra is not installed.
    """
    fitz_mod = require_dependency("fitz", "pdf")

    if use_layout_engine:
        try:
            import pymupdf4llm

            return pymupdf4llm.to_markdown(str(file_path))
        except ImportError:
            warnings.warn(
                "pymupdf4llm is not installed; falling back to the built-in PDF engine. "
                "Install it with: pip install any-to-markdown[pdf]",
                UserWarning,
                stacklevel=2,
            )

    parts: List[str] = []
    path = Path(file_path)
    file_name = path.name
    with fitz_mod.open(path) as doc:
        for page_no, page in enumerate(doc, start=1):
            parts.append(f"---\nsource: {file_name}\npage: {page_no}\ntype: pdf\n---\n\n")

            # 1. Spatial Table Detection
            tabs = page.find_tables()
            table_bboxes = [t.bbox for t in tabs.tables]

            # 2. Extract Text while excluding table areas to prevent duplicate data
            items = _get_pdf_items(page, ignore_bboxes=table_bboxes)

            # 3. Add Table data to the sortable items list
            for tab in tabs.tables:
                df = tab.to_pandas()
                if not df.empty:
                    # Use the top Y-coordinate for correct interleaving
                    items.append((tab.bbox[1], f"\n\n{df.to_markdown(index=False)}\n\n"))

            # 4. Global sort to reconstruct the page layout accurately
            items.sort(key=lambda x: x[0])
            text = "".join(item[1] for item in items)

            if text.strip():
                parts.append(text + "\n\n")

            # 5. Visual/OCR Fallback: Only trigger when text extraction
            #    yielded very little content, indicating the page is likely
            #    image-based / scanned. The previous condition also fired when
            #    *any* image was present (logos, charts, decorations), which
            #    caused extreme slowness on image-rich but text-normal PDFs.
            if len(text.strip()) < 100:
                with tempfile.NamedTemporaryFile(suffix=f"_page_{page_no}.png", delete=False) as temp_image:
                    img_path = Path(temp_image.name)

                try:
                    # Render the page at 2x zoom for high-quality OCR
                    pix = page.get_pixmap(matrix=fitz_mod.Matrix(2, 2))
                    pix.save(img_path)
                    ocr_text = handle_image(img_path)
                    if ocr_text and ocr_text not in text:
                        parts.append(f"\n> [Visual Content OCR]:\n> {ocr_text}\n\n")
                except Exception as exc:
                    # Degrade gracefully: any OCR failure (missing dependency,
                    # tesseract binary not found, processing error, etc.)
                    # should never abort the entire PDF conversion.
                    warnings.warn(f"Skipping OCR fallback for page {page_no}: {exc}", UserWarning, stacklevel=2)
                finally:
                    img_path.unlink(missing_ok=True)
    return "".join(parts)


def handle_notebook(file_path: str | Path) -> str:
    """Parses Jupyter Notebooks (.ipynb) and converts cells to Markdown.

    Args:
        file_path: Path to the .ipynb file.

    Returns:
        Consolidated Markdown content.
    """
    import nbformat

    path = Path(file_path)
    with path.open("r", encoding="utf-8") as f:
        nb = nbformat.read(f, as_version=4)

    parts: List[str] = []
    for cell in nb.cells:
        if cell.cell_type == "markdown":
            parts.append(cell.source + "\n\n")
        elif cell.cell_type == "code":
            # Default to python if language not specified in metadata
            lang: str = nb.metadata.get("language_info", {}).get("name", "python")
            parts.append(f"```{lang}\n{cell.source}\n```\n\n")
    return "".join(parts)


def handle_code(file_path: str | Path) -> str:
    """Reads source code files and wraps them in Markdown code blocks.

    Args:
        file_path: Path to the source file.

    Returns:
        The content wrapped in a fenced code block with language detection.
    """
    ext_to_lang = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".cpp": "cpp",
        ".c": "c",
        ".h": "c",
        ".hpp": "cpp",
        ".rs": "rust",
        ".go": "go",
        ".java": "java",
        ".rb": "ruby",
        ".php": "php",
        ".sh": "bash",
        ".sql": "sql",
        ".yaml": "yaml",
        ".yml": "yaml",
        ".json": "json",
        ".xml": "xml",
        ".html": "html",
        ".css": "css",
        ".md": "markdown",
    }
    path = Path(file_path)
    ext = path.suffix.lower()
    lang = ext_to_lang.get(ext, ext.lstrip(".") or "text")

    with path.open("r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    return f"```{lang}\n{content}\n```\n\n"


class _MarkdownHTMLParser(HTMLParser):
    """Converts a practical subset of HTML into Markdown.

    Supported: headings, paragraphs, line breaks, bold/italic, inline code,
    pre blocks, blockquotes, links, nested ordered/unordered lists, tables,
    and horizontal rules. Non-content tags (script, style, head, ...) are
    stripped entirely. Unknown tags degrade gracefully to their text content.
    """

    _SKIPPED_TAGS = {"script", "style", "head", "title", "meta", "link", "noscript"}
    _HEADING_PREFIXES = {"h1": "#", "h2": "##", "h3": "###", "h4": "####", "h5": "#####", "h6": "######"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: List[str] = []
        self._skip_depth = 0
        self._in_pre = False
        self._list_stack: List[str] = []
        self._ol_counters: List[int] = []
        self._href: Optional[str] = None
        self._link_text: List[str] = []
        self._table_rows: Optional[List[List[str]]] = None
        self._current_row: Optional[List[str]] = None
        self._current_cell: Optional[List[str]] = None

    def get_markdown(self) -> str:
        """Returns the accumulated Markdown with normalized blank lines."""
        text = "".join(self._parts)
        return re.sub(r"\n{3,}", "\n\n", text).strip()

    def _emit(self, text: str) -> None:
        """Routes text to the innermost active sink (link, table cell, or body)."""
        if self._href is not None:
            self._link_text.append(text)
        elif self._current_cell is not None:
            self._current_cell.append(text)
        else:
            self._parts.append(text)

    def _emit_block(self, text: str) -> None:
        """Emits text to the current cell or the body, bypassing the link sink."""
        if self._current_cell is not None:
            self._current_cell.append(text)
        else:
            self._parts.append(text)

    def handle_starttag(self, tag: str, attrs: List[Tuple[str, Optional[str]]]) -> None:
        if tag in self._SKIPPED_TAGS:
            self._skip_depth += 1
            return
        if self._skip_depth:
            return

        if tag in self._HEADING_PREFIXES:
            self._emit(f"\n\n{self._HEADING_PREFIXES[tag]} ")
        elif tag == "p":
            self._emit("\n\n")
        elif tag == "br":
            self._emit("  \n")
        elif tag == "hr":
            self._emit("\n\n---\n\n")
        elif tag in ("strong", "b"):
            self._emit("**")
        elif tag in ("em", "i"):
            self._emit("*")
        elif tag == "code" and not self._in_pre:
            self._emit("`")
        elif tag == "pre":
            self._in_pre = True
            self._emit("\n\n```\n")
        elif tag == "blockquote":
            self._emit("\n\n> ")
        elif tag in ("ul", "ol"):
            self._list_stack.append(tag)
            self._ol_counters.append(0)
        elif tag == "li":
            indent = "  " * max(len(self._list_stack) - 1, 0)
            if self._list_stack and self._list_stack[-1] == "ol":
                self._ol_counters[-1] += 1
                self._emit(f"\n{indent}{self._ol_counters[-1]}. ")
            else:
                self._emit(f"\n{indent}- ")
        elif tag == "a":
            self._href = next((value or "" for name, value in attrs if name == "href"), "")
            self._link_text = []
        elif tag == "table":
            self._table_rows = []
        elif tag == "tr" and self._table_rows is not None:
            self._current_row = []
        elif tag in ("td", "th") and self._current_row is not None:
            self._current_cell = []

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIPPED_TAGS:
            self._skip_depth = max(self._skip_depth - 1, 0)
            return
        if self._skip_depth:
            return

        if tag in self._HEADING_PREFIXES or tag in ("p", "blockquote"):
            self._emit("\n\n")
        elif tag in ("strong", "b"):
            self._emit("**")
        elif tag in ("em", "i"):
            self._emit("*")
        elif tag == "code" and not self._in_pre:
            self._emit("`")
        elif tag == "pre":
            self._in_pre = False
            self._emit("\n```\n\n")
        elif tag in ("ul", "ol"):
            if self._list_stack:
                self._list_stack.pop()
                self._ol_counters.pop()
            self._emit("\n\n")
        elif tag == "a" and self._href is not None:
            text = "".join(self._link_text).strip()
            href = self._href
            self._href = None
            self._emit_block(f"[{text}]({href})" if text else "")
        elif tag in ("td", "th") and self._current_cell is not None and self._current_row is not None:
            self._current_row.append("".join(self._current_cell).strip().replace("|", "\\|"))
            self._current_cell = None
        elif tag == "tr" and self._current_row is not None and self._table_rows is not None:
            self._table_rows.append(self._current_row)
            self._current_row = None
        elif tag == "table":
            self._flush_table()

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        if self._in_pre:
            self._emit(data)
            return
        collapsed = re.sub(r"[ \t\r\n]+", " ", data)
        if collapsed.strip():
            self._emit(collapsed)

    def _flush_table(self) -> None:
        rows = [row for row in (self._table_rows or []) if row]
        self._table_rows = None
        if not rows:
            return
        width = max(len(row) for row in rows)
        padded = [row + [""] * (width - len(row)) for row in rows]
        lines = [
            "| " + " | ".join(padded[0]) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]
        lines.extend("| " + " | ".join(row) + " |" for row in padded[1:])
        self._emit_block("\n\n" + "\n".join(lines) + "\n\n")


def handle_html(file_path: str | Path) -> str:
    """Converts HTML files into structured Markdown.

    Uses a stdlib html.parser-based converter (no extra dependency) that maps
    headings, paragraphs, emphasis, links, lists, tables, code/pre blocks, and
    blockquotes to their Markdown equivalents. Script and style content is
    stripped. Decoding follows the same tolerant strategy as handle_text.

    Args:
        file_path: Path to the .html/.htm file.

    Returns:
        The converted Markdown content.
    """
    raw = Path(file_path).read_bytes()
    try:
        html = raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        html = raw.decode("latin-1")

    parser = _MarkdownHTMLParser()
    parser.feed(html)
    parser.close()
    return f"{parser.get_markdown()}\n\n"


def handle_audio(file_path: str | Path, whisper_model: Optional[str] = None) -> str:
    """Transcribes audio files using the local Faster-Whisper engine.

    Args:
        file_path: Path to the audio file.
        whisper_model: Optional Whisper model size (e.g. 'tiny', 'small',
            'medium'). See get_whisper_model for default resolution.

    Returns:
        The generated transcript text.

    Raises:
        MissingDependencyError: If the 'audio' extra is not installed.
    """
    model = get_whisper_model(whisper_model)
    # segments is a generator; we consume it to get the full transcript
    segments, _ = model.transcribe(str(file_path), vad_filter=True)
    return "".join(segment.text + "\n" for segment in segments)


def handle_video(file_path: str | Path, whisper_model: Optional[str] = None) -> str:
    """Processes video files by extracting the audio stream and transcribing it.

    Args:
        file_path: Path to the video file.
        whisper_model: Optional Whisper model size (see get_whisper_model).

    Returns:
        The generated transcript text extracted from the video's audio.

    Raises:
        MissingDependencyError: If FFmpeg is not installed.
        RuntimeError: If FFmpeg fails, including its stderr output for
            diagnostics.
    """
    audio_path = Path(tempfile.gettempdir()) / f"temp_{uuid4()}.mp3"

    try:
        # Extract the audio track without re-encoding video to save time.
        # stderr is captured (not discarded) so failures are diagnosable.
        try:
            proc = subprocess.run(
                [
                    "ffmpeg",
                    "-i",
                    str(file_path),
                    "-q:a",
                    "0",
                    "-map",
                    "a",
                    str(audio_path),
                    "-y",
                ],
                capture_output=True,
                text=True,
            )
        except FileNotFoundError as e:
            raise MissingDependencyError(
                "FFmpeg is required to process video files. "
                "Install it from https://ffmpeg.org and ensure it is on your PATH."
            ) from e

        if proc.returncode != 0:
            stderr_tail = (proc.stderr or "").strip()[-500:]
            raise RuntimeError(f"ffmpeg failed with exit code {proc.returncode}: {stderr_tail}")

        return handle_audio(audio_path, whisper_model)
    finally:
        # Guarantee cleanup of the heavy audio file
        audio_path.unlink(missing_ok=True)


def extract_youtube_id(url: str) -> Optional[str]:
    """Extracts the 11-character YouTube video ID using a robust regex.

    Args:
        url: The YouTube URL or potential ID.

    Returns:
        The extracted 11-char ID, or None if no match found.
    """
    regex = r"(?:https?:\/\/)?(?:www\.)?(?:youtube\.com\/(?:[^\/\n\s]+\/\S+\/|(?:v|e(?:mbed)?)\/|\S*?[?&]v=)|youtu\.be\/)([a-zA-Z0-9_-]{11})"
    match = re.search(regex, url)
    return match.group(1) if match else None


def handle_youtube(video_id_or_url: str) -> str:
    """Fetches transcripts for YouTube videos via the internal API.

    Args:
        video_id_or_url: A YouTube URL or 11-character video ID.

    Returns:
        The transcript text fetched from YouTube.

    Raises:
        MissingDependencyError: If the 'youtube' extra is not installed.
        TranscriptUnavailableError: If the transcript cannot be retrieved.
            The original exception is preserved as __cause__.
    """
    yta = require_dependency("youtube_transcript_api", "youtube")
    no_transcript_found = getattr(yta, "NoTranscriptFound", Exception)

    video_id = extract_youtube_id(video_id_or_url) or video_id_or_url
    api = yta.YouTubeTranscriptApi()

    try:
        transcript_list = api.list(video_id)

        try:
            # Priority 1: English
            transcript = transcript_list.find_transcript(["en"])
        except no_transcript_found:
            # Priority 2: Any available language
            transcript = next(iter(transcript_list))

        # youtube-transcript-api >= 1.x returns FetchedTranscriptSnippet objects,
        # which expose text via an attribute rather than dict access.
        data = transcript.fetch()
        text = " ".join(chunk.text for chunk in data)
        return f"\n\n{text}\n\n"

    except Exception as e:
        # Typed error with the original exception chained (no information loss).
        raise TranscriptUnavailableError(
            f"No transcript available for YouTube video {video_id}: {e}. "
            "Use handle_yt_local() to transcribe it locally instead."
        ) from e
